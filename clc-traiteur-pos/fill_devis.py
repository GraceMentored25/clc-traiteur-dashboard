# -*- coding: utf-8 -*-
"""
fill_devis.py
Usage: python fill_devis.py <devis_json_path> <output_pptx_path>

Lit Devis_modele.pptx, conserve uniquement les slides nécessaires au devis,
duplique les sections manquantes si le devis en a plus de 3,
et remplit toutes les données.
"""

import sys, json, os, copy
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.oxml.ns import qn
from lxml import etree

TEMPLATE = os.path.join(os.path.dirname(__file__), "public", "Devis_modele.pptx")

# ── Formatage ──────────────────────────────────────────────────────────────
MOIS = ["", "janvier", "février", "mars", "avril", "mai", "juin",
        "juillet", "août", "septembre", "octobre", "novembre", "décembre"]

def fmt_money(amount):
    return f"{int(round(amount)):,} €".replace(",", " ")

def fmt_date(iso):
    try:
        p = iso.split("-")
        return f"{int(p[2])} {MOIS[int(p[1])]} {p[0]}"
    except Exception:
        return iso

# ── Helpers shapes ─────────────────────────────────────────────────────────
def set_shape_text(shape, text):
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    for para in tf.paragraphs[1:]:
        p = para._p
        p.getparent().remove(p)
    para = tf.paragraphs[0]
    runs = para.runs
    if runs:
        for run in runs[1:]:
            run._r.getparent().remove(run._r)
        runs[0].text = text
    else:
        r = para._p
        a_r = etree.SubElement(r, qn("a:r"))
        a_t = etree.SubElement(a_r, qn("a:t"))
        a_t.text = text

def find_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None

def get_text(slide, name):
    s = find_shape(slide, name)
    return s.text_frame.text if s and s.has_text_frame else ""

def slide_has(slide, kw):
    for s in slide.shapes:
        if s.has_text_frame and kw.lower() in s.text_frame.text.lower():
            return True
    return False

# ── Regroupement items par section ─────────────────────────────────────────
def group_sections(devis):
    from collections import OrderedDict
    sec_map = OrderedDict()
    for item in devis.get("items", []):
        label = item.get("section") or "Prestation"
        if label not in sec_map:
            sec_map[label] = []
        sec_map[label].append(item)
    sections = []
    for label, items in sec_map.items():
        subtotal = sum(i["quantity"] * i["unitPrice"] for i in items)
        sections.append({"label": label, "items": items, "subtotal": subtotal})
    return sections

# ── Duplication XML d'une slide ────────────────────────────────────────────
def duplicate_event_slide(prs, source_slide):
    """Duplique une slide événement et l'insère juste après elle."""
    template = source_slide._element
    new_el = copy.deepcopy(template)

    # Ajouter la nouvelle slide dans le package
    from pptx.opc.part import Part
    from pptx.opc.packuri import PackURI
    import random, string
    rnd = "".join(random.choices(string.ascii_lowercase, k=6))
    new_partname = f"/ppt/slides/slideCopy_{rnd}.xml"

    # Utiliser la même présentation layout
    new_slide = prs.slides._sldIdLst

    # Méthode : copier via XML directement
    src_idx = list(prs.slides).index(source_slide)

    # Trouver le rId de la source
    src_rId = prs.slides._sldIdLst[src_idx].get("r:id")
    src_part = prs.part.related_parts[src_rId]

    # Créer une nouvelle part avec le même XML
    from pptx.opc.part import XmlPart
    new_part = copy.deepcopy(src_part)
    # Enregistrer dans le package avec un nouveau nom
    new_rId = prs.part.relate_to(new_part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide")

    # Ajouter l'entrée dans sldIdLst après la source
    max_id = max(int(el.get("id")) for el in prs.slides._sldIdLst) + 1
    new_sldId = etree.SubElement(prs.slides._sldIdLst, qn("p:sldId"))
    new_sldId.set("id", str(max_id))
    new_sldId.set(qn("r:id"), new_rId)

    # Déplacer après la source
    prs.slides._sldIdLst.remove(new_sldId)
    prs.slides._sldIdLst.insert(src_idx + 1, new_sldId)

    return list(prs.slides)[src_idx + 1]

# ── Remplissage Cover ──────────────────────────────────────────────────────
def fill_cover(slide, devis):
    mapping = {
        "Rectangle 10": devis["clientName"],
        "Rectangle 13": fmt_date(devis["eventDate"]),
        "Rectangle 16": devis["eventType"],
        "Rectangle 19": devis.get("lieu", "France"),
        "Rectangle 22": devis.get("clientPhone", ""),
        "Rectangle 25": devis["id"],
    }
    for name, val in mapping.items():
        s = find_shape(slide, name)
        if s:
            set_shape_text(s, val)

# ── Remplissage slide événement (1 section) ────────────────────────────────
# Le template a 3 sections par slide événement. On remplit section par section.
# Les sections supplémentaires au-delà de 3 génèrent une slide dupliquée.

SECTION_SHAPES = [
    # (titre, desc, subtotal, [(plat1, conv1), (plat2, conv2), ...])
    ("Rectangle 16", "Rectangle 17", "Rectangle 20",
     [("Rectangle 22","Rectangle 23"),("Rectangle 25","Rectangle 26"),
      ("Rectangle 28","Rectangle 29"),("Rectangle 31","Rectangle 32")]),
    ("Rectangle 36", "Rectangle 37", "Rectangle 40",
     [("Rectangle 42","Rectangle 43"),("Rectangle 45","Rectangle 46"),
      ("Rectangle 48","Rectangle 49"),("Rectangle 51","Rectangle 52")]),
    ("Rectangle 56", "Rectangle 57", "Rectangle 60",
     [("Rectangle 62","Rectangle 63"),("Rectangle 65","Rectangle 66"),
      ("Rectangle 68","Rectangle 69")]),
]

def fill_event_slide_chunk(slide, devis, chunk_sections, chunk_offset=0):
    """Remplit une slide événement avec jusqu'à 3 sections depuis chunk_sections."""
    # Header
    for name, val in [
        ("Rectangle 4",  devis["eventType"].upper()),
        ("Rectangle 6",  f"{devis['guestCount']} convives"),
        ("Rectangle 9",  fmt_date(devis["eventDate"])),
        ("Rectangle 12", devis.get("lieu", "France")),
    ]:
        s = find_shape(slide, name)
        if s:
            set_shape_text(s, val)

    total_chunk = 0
    for slot_idx, (st, sd, ss, plats) in enumerate(SECTION_SHAPES):
        if slot_idx < len(chunk_sections):
            sec = chunk_sections[slot_idx]
            total_chunk += sec["subtotal"]
            abs_idx = chunk_offset + slot_idx + 1  # numéro de section global

            s = find_shape(slide, st)
            if s:
                set_shape_text(s, f"{abs_idx}. {sec['label']}")
            s = find_shape(slide, sd)
            if s:
                set_shape_text(s, sec.get("desc", ""))
            s = find_shape(slide, ss)
            if s:
                set_shape_text(s, fmt_money(sec["subtotal"]))

            items = sec["items"]
            for pi, (pn, cn) in enumerate(plats):
                ps = find_shape(slide, pn)
                cs = find_shape(slide, cn)
                if pi < len(items):
                    item = items[pi]
                    if ps:
                        set_shape_text(ps, item["dishName"])
                    if cs:
                        set_shape_text(cs, f"{item['quantity']} convives")
                else:
                    if ps:
                        set_shape_text(ps, "")
                    if cs:
                        set_shape_text(cs, "")
        else:
            # Vider ce slot
            for n in [st, sd, ss]:
                s = find_shape(slide, n)
                if s:
                    set_shape_text(s, "")
            for pn, cn in plats:
                ps = find_shape(slide, pn)
                cs = find_shape(slide, cn)
                if ps:
                    set_shape_text(ps, "")
                if cs:
                    set_shape_text(cs, "")

    # Sous-total slide
    for name in ["Rectangle 71", "Rectangle 74", "Rectangle 72", "Rectangle 75"]:
        s = find_shape(slide, name)
        if s:
            t = s.text_frame.text if s.has_text_frame else ""
            if "SOUS-TOTAL" in t.upper():
                continue
            if "€" in t or t.strip() == "" or t.replace(" ", "").replace(" ", "").replace("€", "").isdigit():
                set_shape_text(s, fmt_money(total_chunk))

# ── Remplissage récapitulatif ──────────────────────────────────────────────
def fill_recap(slide, devis, sections):
    total_event = sum(s["subtotal"] for s in sections)

    for name, val in [
        ("Rectangle 5",  devis["eventType"].upper()),
        ("Rectangle 7",  f"{devis['guestCount']} convives"),
        ("Rectangle 9",  fmt_date(devis["eventDate"])),
        ("Rectangle 11", devis.get("lieu", "France")),
    ]:
        s = find_shape(slide, name)
        if s:
            set_shape_text(s, val)

    rows = [
        ("Rectangle 15","Rectangle 16","Rectangle 17","Rectangle 18"),
        ("Rectangle 21","Rectangle 22","Rectangle 23","Rectangle 24"),
        ("Rectangle 27","Rectangle 28","Rectangle 29","Rectangle 30"),
    ]
    for ri, (rn, rl, rd, rp) in enumerate(rows):
        if ri < len(sections):
            sec = sections[ri]
            for n, v in [(rn, str(ri+1)), (rl, sec["label"]),
                         (rd, sec.get("desc", "")), (rp, fmt_money(sec["subtotal"]))]:
                s = find_shape(slide, n)
                if s:
                    set_shape_text(s, v)
        else:
            for n in [rn, rl, rd, rp]:
                s = find_shape(slide, n)
                if s:
                    set_shape_text(s, "")

    # Si > 3 sections, afficher la somme dans le sous-total événement
    s = find_shape(slide, "Rectangle 33")
    if s:
        set_shape_text(s, fmt_money(total_event))

    s = find_shape(slide, "Rectangle 55")
    if s:
        set_shape_text(s, fmt_money(devis["totalTTC"]))

# ── Remplissage acompte ────────────────────────────────────────────────────
def fill_acompte(slide, devis, sections):
    total_event = sum(s["subtotal"] for s in sections)
    total_ttc   = devis["totalTTC"]
    a30 = round(total_ttc * 0.30)
    a40 = round(total_ttc * 0.40)

    for name, val in [
        ("Rectangle 5",  devis["eventType"].upper()),
        ("Rectangle 8",  fmt_money(total_ttc)),
        ("Rectangle 9",  fmt_money(total_event) + "  événement"),
        ("Rectangle 14", fmt_money(a30)),
        ("Rectangle 26", fmt_money(a30)),
        ("Rectangle 32", fmt_money(a40)),
        ("Rectangle 38", fmt_money(a30)),
    ]:
        s = find_shape(slide, name)
        if s:
            set_shape_text(s, val)

# ── Slides du template ─────────────────────────────────────────────────────
EVENT_SLIDES = {
    "mariage":           (1, 7, 8),
    "anniversaire":      (2, 9, 10),
    "baby shower":       (3, 11, 12),
    "séminaire":         (4, 13, 14),
    "seminaire":         (4, 13, 14),
    "réception privée":  (5, 15, 16),
    "reception privee":  (5, 15, 16),
}

def match_event(event_type):
    et = event_type.lower()
    for key, val in EVENT_SLIDES.items():
        if key in et or et in key:
            return val
    return (1, 7, 8)  # défaut Mariage

# ── Suppression d'une slide ────────────────────────────────────────────────
def delete_slide_at(prs, idx):
    try:
        rId = prs.slides._sldIdLst[idx].get("r:id")
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[idx])
        if rId in prs.part.related_parts:
            del prs.part.related_parts[rId]
    except Exception as e:
        print(f"  [warn] delete slide {idx}: {e}")


# ── Renumérotation des pages ───────────────────────────────────────────────
def renumber_pages(prs):
    """Met à jour le numéro de page global dans chaque slide.
    Chaque slide a une Ellipse dont le texte est le numéro de page.
    On cherche l'Ellipse avec la valeur la plus grande (= numéro global).
    """
    for page_num, slide in enumerate(prs.slides, start=1):
        # Trouver toutes les ellipses avec un contenu numérique
        ellipses = []
        for shape in slide.shapes:
            if 'Ellipse' in shape.name and shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t.isdigit():
                    ellipses.append((int(t), shape))
        if ellipses:
            # L'ellipse avec le plus grand numéro = numéro de page global
            ellipses.sort(key=lambda x: x[0])
            page_shape = ellipses[-1][1]
            set_shape_text(page_shape, str(page_num))

# ── Main ───────────────────────────────────────────────────────────────────
def main():
    if len(sys.argv) < 3:
        print("Usage: python fill_devis.py <devis.json> <output.pptx>")
        sys.exit(1)

    with open(sys.argv[1], encoding="utf-8") as f:
        devis = json.load(f)

    sections   = group_sections(devis)
    event_type = devis.get("eventType", "Mariage")
    ev_idx, rc_idx, ac_idx = match_event(event_type)

    prs    = Presentation(TEMPLATE)
    total  = len(list(prs.slides))

    # Slides à garder (indices du template)
    SLIDE_COVER      = 0
    SLIDE_PRESTATION = 6
    SLIDE_MENTIONS   = 17
    SLIDE_SIGNATURE  = 18
    SLIDE_LEGENDE    = 19

    keep = {SLIDE_COVER, ev_idx, rc_idx, ac_idx,
            SLIDE_PRESTATION, SLIDE_MENTIONS, SLIDE_SIGNATURE}

    # Supprimer dans l'ordre inverse
    to_del = sorted([i for i in range(total) if i not in keep], reverse=True)
    for idx in to_del:
        delete_slide_at(prs, idx)

    # ── Duplication si >3 sections ─────────────────────────────────────────
    # Après suppression, il reste : cover, event, prestation, recap, acompte, mentions, signature
    # On identifie la slide événement et on la duplique si nécessaire
    event_slide_obj = None
    for sl in prs.slides:
        if (slide_has(sl, event_type.upper()[:6]) or
                (not slide_has(sl, "DEVIS") and
                 not slide_has(sl, "RÉCAP") and
                 not slide_has(sl, "ACOMPTE") and
                 not slide_has(sl, "MENTIONS") and
                 not slide_has(sl, "ACCORD") and
                 not slide_has(sl, "PRESTATION"))):
            event_slide_obj = sl
            break

    # Créer les slides événement nécessaires (une slide = 3 sections max)
    n_chunks = max(1, (len(sections) + 2) // 3)
    event_slides = [event_slide_obj]

    if n_chunks > 1 and event_slide_obj:
        for _ in range(n_chunks - 1):
            new_sl = duplicate_event_slide(prs, event_slides[-1])
            event_slides.append(new_sl)

    # ── Remplissage ────────────────────────────────────────────────────────
    for sl in prs.slides:
        if slide_has(sl, "Rectangle 10") and slide_has(sl, "DEVIS"):
            fill_cover(sl, devis)
        elif slide_has(sl, "RÉCAPITULATIF") or slide_has(sl, "RECAPITULATIF"):
            fill_recap(sl, devis, sections)
        elif slide_has(sl, "ACOMPTE") and slide_has(sl, "ÉCHÉANCIER"):
            fill_acompte(sl, devis, sections)
        elif slide_has(sl, "MENTIONS") or slide_has(sl, "LÉGALES"):
            pass  # statique
        elif slide_has(sl, "ACCORD") or slide_has(sl, "SIGNATURE"):
            pass  # statique
        elif slide_has(sl, "PRESTATION"):
            pass  # statique
        else:
            # Slide(s) événement
            try:
                chunk_idx = event_slides.index(sl)
                chunk = sections[chunk_idx * 3 : chunk_idx * 3 + 3]
                fill_event_slide_chunk(sl, devis, chunk, chunk_offset=chunk_idx * 3)
            except ValueError:
                pass

    renumber_pages(prs)
    prs.save(sys.argv[2])
    print(f"OK:{sys.argv[2]}")

if __name__ == "__main__":
    main()
